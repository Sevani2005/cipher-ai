import os
import fitz  # PyMuPDF
import cv2
import json
import re
import math
import traceback
import io
from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from docx import Document
import google.generativeai as genai

try:
    from groq import Groq
except ImportError:
    Groq = None

# Helper to format sensitive text into 'X' characters, keeping format
def format_as_xxxx(text, category):
    result = []
    for char in text:
        if char.isdigit():
            result.append("X")
        elif char.isalpha():
            result.append("X")
        else:
            result.append(char)
    return "".join(result)

def get_text_replacement(text, category, mask_style):
    if mask_style == "blur":
        return f"[BLURRED_{category}]"
    elif mask_style == "redact":
        return f"[REDACTED_{category}]"
    elif mask_style == "xxxx":
        return format_as_xxxx(text, category)
    return f"[{category}]"

def is_category_enabled(category, mask_options):
    cat = category.upper()
    if cat == "NAME":
        return mask_options.get("name", True)
    elif cat == "EMAIL":
        return mask_options.get("email", True)
    elif cat in ["PHONE", "TELEPHONE"]:
        return mask_options.get("phone", True)
    elif cat in ["CREDIT_CARD", "BANK", "CVV", "FINANCIAL", "AADHAAR_SSN", "SSN", "GOVT_ID", "AADHAAR", "PAN", "PASSPORT"]:
        return mask_options.get("financial", True)
    elif cat in ["API_KEY", "PASSWORD", "CREDENTIALS", "SECRET"]:
        return mask_options.get("credentials", True)
    elif cat in ["ADDRESS", "LOCATION"]:
        return mask_options.get("address", True)
    else:
        return mask_options.get("other", True)

class MaskingAgent:
    def __init__(self, api_key=None, api_provider="local"):
        self.api_provider = api_provider or "local"
        
        # Fallback to server-side environment variables if no key is provided in browser
        if not api_key or api_key.strip() in ["", "null", "undefined"]:
            if self.api_provider == "groq":
                self.api_key = os.getenv("GROQ_API_KEY")
            elif self.api_provider == "gemini":
                self.api_key = os.getenv("GEMINI_API_KEY") or os.getenv("GOOGLE_API_KEY")
            else:
                self.api_key = None
        else:
            self.api_key = api_key

        self.has_gemini = False
        self.has_groq = False

        if self.api_key and self.api_key.strip() not in ["", "null", "undefined"]:
            if self.api_provider == "gemini":
                try:
                    genai.configure(api_key=self.api_key)
                    self.has_gemini = True
                except Exception:
                    self.has_gemini = False
            elif self.api_provider == "groq" and Groq:
                try:
                    self.groq_client = Groq(api_key=self.api_key)
                    self.has_groq = True
                except Exception:
                    self.has_groq = False

    def set_api_key(self, api_key, api_provider):
        self.api_provider = api_provider or "local"
        if api_key and api_key.strip() not in ["", "null", "undefined"]:
            self.api_key = api_key
            if self.api_provider == "gemini":
                try:
                    genai.configure(api_key=api_key)
                    self.has_gemini = True
                    self.has_groq = False
                except Exception:
                    self.has_gemini = False
            elif self.api_provider == "groq" and Groq:
                try:
                    self.groq_client = Groq(api_key=api_key)
                    self.has_groq = True
                    self.has_gemini = False
                except Exception:
                    self.has_groq = False
        else:
            self.api_key = None
            self.has_gemini = False
            self.has_groq = False

    def _get_gemini_model(self):
        if not self.has_gemini:
            raise ValueError("Gemini API key is not configured.")
        return genai.GenerativeModel('gemini-2.5-flash')

    def analyze_text_for_pii(self, text, custom_keywords=None, custom_regex=None, mask_options=None, user_prompt=None):
        """
        Routes the text to the selected active engine: Groq LLM, Gemini AI, or Local Regex fallback.
        """
        if self.api_provider == "groq" and self.has_groq:
            return self.analyze_text_for_pii_with_groq(text, custom_keywords, custom_regex, mask_options, user_prompt)
        elif self.api_provider == "gemini" and self.has_gemini:
            return self.analyze_text_for_pii_with_gemini(text, custom_keywords, custom_regex, mask_options, user_prompt)
        else:
            return self.local_analyze_text_for_pii(text, custom_keywords, custom_regex, mask_options)

    def analyze_text_for_pii_with_groq(self, text, custom_keywords=None, custom_regex=None, mask_options=None, user_prompt=None):
        """
        Analyze text using Groq Llama-3.1-8b-instant model. Extremely fast execution!
        """
        logs = []
        if not text.strip():
            return {"sensitive_elements": []}, logs

        logs.append({"type": "info", "message": "[Agent] Activating Groq Cloud LLM Engine (llama-3.1-8b-instant)..."})
        
        prompt = f"""
You are an advanced security AI. Analyze the following text and identify all sensitive or personally identifiable information (PII).
Categories to detect:
1. NAME (e.g. Person names)
2. EMAIL (e.g. Email addresses)
3. PHONE (e.g. Phone numbers, mobile numbers)
4. CREDIT_CARD (e.g. Credit card numbers, banking details)
5. AADHAAR_SSN (e.g. SSN numbers, Aadhaar numbers, PAN, government IDs)
6. ADDRESS (e.g. Physical addresses, specific locations)
7. API_KEY (e.g. API keys, secrets, private tokens)
8. PASSWORD (e.g. Password text)
9. OTHER (e.g. Dates of birth, metadata, or other personal details)

Return a valid JSON object with a single key "sensitive_elements" which contains a list of objects.
Each object must have:
- "text": The EXACT substring as it appears in the source text. It must match the source character-for-character, including casing.
- "category": One of the category strings above (e.g., "NAME", "EMAIL", "PHONE", "CREDIT_CARD", "AADHAAR_SSN", "ADDRESS", "API_KEY", "PASSWORD", "OTHER").

Do not include any explanation or markdown formatting, just the raw JSON.

Input text:
{text}
"""
        if user_prompt and user_prompt.strip():
            prompt += f"\n\nCRITICAL USER REQUEST: The user has explicitly requested to mask/hide: '{user_prompt}'. You MUST identify any words, phrases, names, entities, or parts of the text matching this request, and include them in the JSON list with category 'OTHER'."

        try:
            # Query Groq using JSON response format mode
            response = self.groq_client.chat.completions.create(
                model="llama-3.1-8b-instant",
                messages=[
                    {"role": "user", "content": prompt}
                ],
                response_format={"type": "json_object"}
            )
            
            content = response.choices[0].message.content
            data = json.loads(content)
            elements = data.get("sensitive_elements", [])
            logs.append({"type": "success", "message": f"Groq Llama-3.1 identified {len(elements)} sensitive element(s) in document text."})
            
            # Incorporate custom keywords
            if custom_keywords:
                for kw in custom_keywords:
                    kw_clean = kw.strip()
                    if kw_clean and kw_clean.lower() in text.lower():
                        matches = re.finditer(re.escape(kw_clean), text, re.IGNORECASE)
                        for match in matches:
                            matched_text = text[match.start():match.end()]
                            if not any(e["text"].lower() == matched_text.lower() for e in elements):
                                elements.append({"text": matched_text, "category": "OTHER"})
                                
            # Incorporate custom regex
            if custom_regex:
                for pattern_str in custom_regex:
                    pattern_clean = pattern_str.strip()
                    if pattern_clean:
                        try:
                            matches = re.finditer(pattern_clean, text)
                            for match in matches:
                                matched_text = match.group()
                                if not any(e["text"].lower() == matched_text.lower() for e in elements):
                                    elements.append({"text": matched_text, "category": "OTHER"})
                        except Exception as e:
                            logs.append({"type": "warning", "message": f"Skipping invalid regex pattern '{pattern_clean}': {str(e)}"})

            return {"sensitive_elements": elements}, logs
        except Exception as e:
            logs.append({"type": "error", "message": f"Groq API call failed: {str(e)}. Falling back to local rules..."})
            traceback.print_exc()
            return self.local_analyze_text_for_pii(text, custom_keywords, custom_regex, mask_options)

    def analyze_text_for_pii_with_gemini(self, text, custom_keywords=None, custom_regex=None, mask_options=None, user_prompt=None):
        """
        Analyze text using Google Gemini 1.5 Flash.
        """
        logs = []
        if not text.strip():
            return {"sensitive_elements": []}, logs

        logs.append({"type": "info", "message": "[Agent] Activating Gemini Generative AI Engine..."})
        
        prompt = f"""
You are an advanced security AI. Analyze the following text and identify all sensitive or personally identifiable information (PII).
Categories to detect:
1. NAME (e.g. Person names)
2. EMAIL (e.g. Email addresses)
3. PHONE (e.g. Phone numbers, mobile numbers)
4. CREDIT_CARD (e.g. Credit card numbers, banking details)
5. AADHAAR_SSN (e.g. SSN numbers, Aadhaar numbers, PAN, government IDs)
6. ADDRESS (e.g. Physical addresses, specific locations)
7. API_KEY (e.g. API keys, secrets, private tokens)
8. PASSWORD (e.g. Password text)
9. OTHER (e.g. Dates of birth, metadata, or other personal details)

Return a valid JSON object with a single key "sensitive_elements" which contains a list of objects.
Each object must have:
- "text": The EXACT substring as it appears in the source text. It must match the source character-for-character, including casing.
- "category": One of the category strings above (e.g., "NAME", "EMAIL", "PHONE", "CREDIT_CARD", "AADHAAR_SSN", "ADDRESS", "API_KEY", "PASSWORD", "OTHER").

Do not include any explanation or markdown formatting, just the raw JSON.

Input text:
{text}
"""
        if user_prompt and user_prompt.strip():
            prompt += f"\n\nCRITICAL USER REQUEST: The user has explicitly requested to mask/hide: '{user_prompt}'. You MUST identify any words, phrases, names, entities, or parts of the text matching this request, and include them in the JSON list with category 'OTHER'."

        try:
            model = self._get_gemini_model()
            response = model.generate_content(
                prompt,
                generation_config={"response_mime_type": "application/json"}
            )
            data = json.loads(response.text)
            elements = data.get("sensitive_elements", [])
            logs.append({"type": "success", "message": f"Gemini detected {len(elements)} sensitive text element(s)."})
            
            # Incorporate custom keywords
            if custom_keywords:
                for kw in custom_keywords:
                    kw_clean = kw.strip()
                    if kw_clean and kw_clean.lower() in text.lower():
                        matches = re.finditer(re.escape(kw_clean), text, re.IGNORECASE)
                        for match in matches:
                            matched_text = text[match.start():match.end()]
                            if not any(e["text"].lower() == matched_text.lower() for e in elements):
                                elements.append({"text": matched_text, "category": "OTHER"})
                                
            # Incorporate custom regex
            if custom_regex:
                for pattern_str in custom_regex:
                    pattern_clean = pattern_str.strip()
                    if pattern_clean:
                        try:
                            matches = re.finditer(pattern_clean, text)
                            for match in matches:
                                matched_text = match.group()
                                if not any(e["text"].lower() == matched_text.lower() for e in elements):
                                    elements.append({"text": matched_text, "category": "OTHER"})
                        except Exception as e:
                            logs.append({"type": "warning", "message": f"Skipping invalid regex pattern '{pattern_clean}': {str(e)}"})

            return {"sensitive_elements": elements}, logs
        except Exception as e:
            logs.append({"type": "error", "message": f"Gemini API call failed: {str(e)}. Falling back to local rules..."})
            traceback.print_exc()
            return self.local_analyze_text_for_pii(text, custom_keywords, custom_regex, mask_options)

    def local_analyze_text_for_pii(self, text, custom_keywords=None, custom_regex=None, mask_options=None):
        """
        Runs entirely locally using high-fidelity regular expressions. No API Key Required!
        """
        elements = []
        logs = [{"type": "agent", "message": "API Key Missing/Inactive. Activating 100% Local Rules-Based Masking Engine..."}]
        
        if not mask_options:
            mask_options = {}

        # High-fidelity regex dictionary
        regex_map = {
            "EMAIL": r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b',
            "PHONE": r'(?:\+?\d{1,3}[-.\s]?)?\(?\d{3}\)?[-.\s]?\d{3}[-.\s]?\d{4}\b|\b[6-9]\d{9}\b|\+91[-.\s]?[6-9]\d{9}\b',
            "CREDIT_CARD": r'\b\d{4}[- ]?\d{4}[- ]?\d{4}[- ]?\d{4}\b|\b(?:\d[ -]*?){13,16}\b',
            "AADHAAR_SSN": r'\b\d{4}[- ]?\d{4}[- ]?\d{4}\b|\b\d{3}-\d{2}-\d{4}\b|\b[A-Z]{5}[0-9]{4}[A-Z]\b',
            "API_KEY": r'\b(?:key|secret|token|password|passwd|passcode|root_password)\b\s*[:=]\s*["\']?([A-Za-z0-9\-_!@#\$%\^&\*\(\)\+]{6,40})["\']?',
        }

        for category, pattern in regex_map.items():
            if category == "EMAIL" and not mask_options.get("email", True): continue
            if category == "PHONE" and not mask_options.get("phone", True): continue
            if category == "CREDIT_CARD" and not mask_options.get("financial", True): continue
            if category == "AADHAAR_SSN" and not mask_options.get("financial", True): continue
            if category == "API_KEY" and not mask_options.get("credentials", True): continue
            
            matches = re.finditer(pattern, text, re.IGNORECASE)
            for m in matches:
                if category == "API_KEY" and len(m.groups()) > 0 and m.group(1):
                    matched_text = m.group(1)
                else:
                    matched_text = m.group(0)
                
                if matched_text and matched_text.strip():
                    if not any(e["text"] == matched_text for e in elements):
                        elements.append({"text": matched_text, "category": category})

        # Process Custom Keywords
        if custom_keywords:
            for kw in custom_keywords:
                kw_clean = kw.strip()
                if kw_clean:
                    matches = re.finditer(re.escape(kw_clean), text, re.IGNORECASE)
                    for m in matches:
                        matched_text = text[m.start():m.end()]
                        if not any(e["text"] == matched_text for e in elements):
                            elements.append({"text": matched_text, "category": "OTHER"})

        # Process Custom Regex
        if custom_regex:
            for pattern_str in custom_regex:
                pattern_clean = pattern_str.strip()
                if pattern_clean:
                    try:
                        matches = re.finditer(pattern_clean, text)
                        for m in matches:
                            matched_text = m.group(0)
                            if not any(e["text"] == matched_text for e in elements):
                                elements.append({"text": matched_text, "category": "OTHER"})
                    except Exception as e:
                        logs.append({"type": "warning", "message": f"Skipping invalid regex '{pattern_clean}': {str(e)}"})

        logs.append({"type": "success", "message": f"Local parser completed. Detected {len(elements)} sensitive element(s) locally."})
        if mask_options.get("name", True):
            logs.append({"type": "warning", "message": "[Notice] Name detection requires Gemini or Groq API. To mask specific names in local mode, type them in 'Custom Keywords' above."})
            
        return {"sensitive_elements": elements}, logs

    def analyze_image_for_pii(self, pil_image, custom_keywords=None, user_prompt=None):
        """
        Send PIL Image to Gemini to identify visual PII and return bounding box coordinates.
        """
        logs = []
        if not self.has_gemini:
            logs.append({"type": "error", "message": "Visual coordinate-based image masking requires the Gemini API Key."})
            raise ValueError("Gemini API Key is required for Image visual coordinate masking. Please switch to Gemini provider.")

        logs.append({"type": "info", "message": f"Sending image to Gemini for visual PII analysis (Image size: {pil_image.size})..."})
        
        prompt = """
You are an expert security AI agent. Analyze the provided document image and detect any sensitive information.
Categories of sensitive information:
1. NAME (Person names)
2. EMAIL (Email addresses)
3. PHONE (Phone/fax numbers)
4. CREDIT_CARD (Credit cards, banks)
5. AADHAAR_SSN (Aadhaar cards, PAN card, SSN, Passport, Government IDs)
6. ADDRESS (Physical addresses, specific locations)
7. API_KEY (API keys, secrets)
8. PASSWORD (Passwords)
9. OTHER (Other sensitive parameters)

For each sensitive item, detect its exact bounding box and content.
Return a valid JSON object with a single key "sensitive_elements" containing a list of objects.
Each object must contain:
- "text": The text content detected.
- "category": One of "NAME", "EMAIL", "PHONE", "CREDIT_CARD", "AADHAAR_SSN", "ADDRESS", "API_KEY", "PASSWORD", "OTHER".
- "box_2d": The bounding box coordinates of the text, formatted as [ymin, xmin, ymax, xmax].
  These coordinates must be normalized on a scale of 0 to 1000 relative to the image height and width.

If there is no sensitive information, return an empty list.
Do not include any explanation or markdown formatting, just the raw JSON.
"""
        if user_prompt and user_prompt.strip():
            prompt += f"\n\nCRITICAL USER REQUEST: The user has explicitly requested to mask/hide: '{user_prompt}'. You MUST identify any visual features, text segments, areas, faces, barcodes, signatures, or shapes matching or representing this request, locate their coordinates, and return them in the list with category 'OTHER'."

        try:
            model = self._get_gemini_model()
            response = model.generate_content(
                [pil_image, prompt],
                generation_config={"response_mime_type": "application/json"}
            )
            data = json.loads(response.text)
            elements = data.get("sensitive_elements", [])
            logs.append({"type": "success", "message": f"Gemini visually detected {len(elements)} sensitive element(s)."})
            return data, logs
        except Exception as e:
            logs.append({"type": "error", "message": f"Error during Gemini image analysis: {str(e)}"})
            traceback.print_exc()
            return {"sensitive_elements": []}, logs

    def mask_image_file(self, input_path, output_path, mask_options, custom_keywords=None, custom_regex=None, user_prompt=None):
        """
        Masks an image file (PNG/JPG) using Gemini visual PII coordinates.
        """
        logs = [{"type": "agent", "message": "Starting image masking process..."}]
        if not self.has_gemini:
            logs.append({"type": "warning", "message": "[Notice] Visual image pixel masking requires the Gemini API key."})
            try:
                # Fallback: Copy original image so dashboard doesn't crash, and show a clear notice!
                img = Image.open(input_path)
                img.save(output_path)
                logs.append({"type": "warning", "message": "Image preserved without masking. To automatically blur sensitive parts of this image, please switch provider to 'Google Gemini' in the header and save a Gemini API Key."})
            except Exception as e:
                logs.append({"type": "error", "message": f"Failed to generate fallback image: {str(e)}"})
            return logs
            
        try:
            image = Image.open(input_path).convert("RGB")
            width, height = image.size
            
            # Send to Gemini
            data, gemini_logs = self.analyze_image_for_pii(image, custom_keywords, user_prompt)
            logs.extend(gemini_logs)
            
            elements = data.get("sensitive_elements", [])
            mask_style = mask_options.get("style", "redact")
            
            draw = ImageDraw.Draw(image)
            masked_count = 0
            
            for elem in elements:
                text = elem.get("text")
                category = elem.get("category", "OTHER")
                box = elem.get("box_2d")
                
                if not box or len(box) != 4:
                    continue
                    
                if not is_category_enabled(category, mask_options):
                    logs.append({"type": "info", "message": f"Skipping {category}: '{text}' (category disabled in settings)"})
                    continue
                
                ymin, xmin, ymax, xmax = box
                x1 = max(0, int((xmin / 1000.0) * width))
                y1 = max(0, int((ymin / 1000.0) * height))
                x2 = min(width, int((xmax / 1000.0) * width))
                y2 = min(height, int((ymax / 1000.0) * height))
                
                if x2 <= x1 or y2 <= y1:
                    continue
                
                logs.append({"type": "success", "message": f"Masking '{text}' [{category}] at [{x1}, {y1}, {x2}, {y2}]"})
                masked_count += 1
                
                if mask_style == "blur":
                    region = image.crop((x1, y1, x2, y2))
                    radius = max(5, int(math.sqrt((x2-x1) * (y2-y1)) / 10))
                    blurred_region = region.filter(ImageFilter.GaussianBlur(radius=radius))
                    image.paste(blurred_region, (x1, y1, x2, y2))
                elif mask_style == "redact":
                    draw.rectangle([x1, y1, x2, y2], fill="#11111b", outline="#f38ba8", width=1)
                elif mask_style == "xxxx":
                    draw.rectangle([x1, y1, x2, y2], fill="#1e1e2e", outline="#89b4fa", width=1)
                    
            image.save(output_path)
            logs.append({"type": "agent", "message": f"Image masking complete. Masked {masked_count} elements. File saved."})
            return logs
        except Exception as e:
            logs.append({"type": "error", "message": f"Failed to mask image: {str(e)}"})
            traceback.print_exc()
            return logs

    def mask_pdf_file(self, input_path, output_path, mask_options, custom_keywords=None, custom_regex=None, user_prompt=None):
        """
        Masks a PDF document.
        """
        logs = [{"type": "agent", "message": "Starting PDF masking process..."}]
        pdf_mode = mask_options.get("pdf_mode", "vector")
        mask_style = mask_options.get("style", "redact")
        
        if pdf_mode == "visual" and not self.has_gemini:
            logs.append({"type": "error", "message": "Visual Page Flattening mode requires the Gemini API key. Switching to Native Vector Redaction..."})
            pdf_mode = "vector"
            
        try:
            doc = fitz.open(input_path)
            total_pages = len(doc)
            logs.append({"type": "info", "message": f"Opened PDF with {total_pages} page(s). Mode: {pdf_mode.upper()}"})
            
            if pdf_mode == "vector":
                # NATIVE VECTOR REDACTION
                total_redacted = 0
                for page_num in range(total_pages):
                    logs.append({"type": "info", "message": f"Processing Page {page_num + 1}/{total_pages}..."})
                    page = doc[page_num]
                    text = page.get_text()
                    
                    if not text.strip():
                        logs.append({"type": "warning", "message": f"No selectable text on Page {page_num + 1}. Scanned PDFs require visual mode with a Gemini API key."})
                        continue
                        
                    data, text_logs = self.analyze_text_for_pii(text, custom_keywords, custom_regex, mask_options, user_prompt)
                    logs.extend(text_logs)
                    
                    elements = data.get("sensitive_elements", [])
                    page_redacted_count = 0
                    
                    for elem in elements:
                        search_text = elem.get("text")
                        category = elem.get("category", "OTHER")
                        
                        if not is_category_enabled(category, mask_options):
                            continue
                            
                        rects = page.search_for(search_text)
                        if not rects:
                            continue
                            
                        logs.append({"type": "success", "message": f"Redacting '{search_text}' [{category}] ({len(rects)} occurrence(s) on Page {page_num + 1})"})
                        page_redacted_count += len(rects)
                        
                        for rect in rects:
                            if mask_style == "blur":
                                page.add_redact_annot(rect, fill=(0.85, 0.85, 0.85))
                            elif mask_style == "redact":
                                page.add_redact_annot(rect, fill=(0.0, 0.0, 0.0))
                            elif mask_style == "xxxx":
                                label = f"[{category}]"
                                page.add_redact_annot(
                                    rect, 
                                    text=label, 
                                    fill=(0.12, 0.12, 0.18), 
                                    text_color=(0.95, 0.55, 0.66)
                                )
                                
                    if page_redacted_count > 0:
                        page.apply_redactions()
                        total_redacted += page_redacted_count
                        
                doc.save(output_path, garbage=4, deflate=True)
                logs.append({"type": "agent", "message": f"PDF native vector redaction complete. Applied {total_redacted} redactions."})
                
            else:
                # FULL VISUAL REDACTION (Convert pages to images, redact images, merge to PDF)
                logs.append({"type": "agent", "message": "Converting PDF pages to high-resolution images for full visual masking..."})
                temp_pdf = fitz.open()
                
                for page_num in range(total_pages):
                    logs.append({"type": "info", "message": f"Rendering and processing Page {page_num + 1}/{total_pages}..."})
                    page = doc[page_num]
                    
                    pix = page.get_pixmap(dpi=150)
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    width, height = img.size
                    
                    data, img_logs = self.analyze_image_for_pii(img, custom_keywords, user_prompt)
                    logs.extend(img_logs)
                    
                    elements = data.get("sensitive_elements", [])
                    draw = ImageDraw.Draw(img)
                    
                    for elem in elements:
                        text = elem.get("text")
                        category = elem.get("category", "OTHER")
                        box = elem.get("box_2d")
                        
                        if not box or len(box) != 4:
                            continue
                            
                        if not is_category_enabled(category, mask_options):
                            continue
                            
                        ymin, xmin, ymax, xmax = box
                        x1 = max(0, int((xmin / 1000.0) * width))
                        y1 = max(0, int((ymin / 1000.0) * height))
                        x2 = min(width, int((xmax / 1000.0) * width))
                        y2 = min(height, int((ymax / 1000.0) * height))
                        
                        if x2 <= x1 or y2 <= y1:
                            continue
                            
                        logs.append({"type": "success", "message": f"Visual Mask: '{text}' [{category}] at Page {page_num + 1}"})
                        
                        if mask_style == "blur":
                            region = img.crop((x1, y1, x2, y2))
                            radius = max(5, int(math.sqrt((x2-x1) * (y2-y1)) / 10))
                            blurred = region.filter(ImageFilter.GaussianBlur(radius=radius))
                            img.paste(blurred, (x1, y1, x2, y2))
                        elif mask_style == "redact":
                            draw.rectangle([x1, y1, x2, y2], fill="#11111b", outline="#f38ba8", width=1)
                        elif mask_style == "xxxx":
                            draw.rectangle([x1, y1, x2, y2], fill="#1e1e2e", outline="#89b4fa", width=1)
                    
                    img_byte_arr = io.BytesIO()
                    img.save(img_byte_arr, format='PDF')
                    img_byte_arr.seek(0)
                    
                    temp_page_doc = fitz.open("pdf", img_byte_arr.read())
                    temp_pdf.insert_pdf(temp_page_doc)
                
                temp_pdf.save(output_path)
                temp_pdf.close()
                logs.append({"type": "agent", "message": "Visual PDF masking and reconstruction complete."})
                
            doc.close()
            return logs
        except Exception as e:
            logs.append({"type": "error", "message": f"Failed to mask PDF: {str(e)}"})
            traceback.print_exc()
            return logs

    def mask_pptx_file(self, input_path, output_path, mask_options, custom_keywords=None, custom_regex=None, user_prompt=None):
        """
        Masks a PowerPoint Presentation (.pptx).
        """
        logs = [{"type": "agent", "message": "Starting PPTX masking process..."}]
        mask_style = mask_options.get("style", "redact")
        
        try:
            prs = Presentation(input_path)
            logs.append({"type": "info", "message": f"Opened PPTX with {len(prs.slides)} slides."})
            
            all_text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            for run in p.runs:
                                if run.text.strip():
                                    all_text_runs.append(run.text)
                                    
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                for p in cell.text_frame.paragraphs:
                                    for run in p.runs:
                                        if run.text.strip():
                                            all_text_runs.append(run.text)
            
            if not all_text_runs:
                logs.append({"type": "warning", "message": "No text detected in presentation shapes."})
                prs.save(output_path)
                return logs
                
            full_text_block = "\n---\n".join(all_text_runs)
            data, gemini_logs = self.analyze_text_for_pii(full_text_block, custom_keywords, custom_regex, mask_options, user_prompt)
            logs.extend(gemini_logs)
            
            sensitive_elements = data.get("sensitive_elements", [])
            if not sensitive_elements:
                logs.append({"type": "info", "message": "No sensitive elements identified in presentation."})
                prs.save(output_path)
                return logs
                
            replace_count = 0
            
            def replace_text_in_run(run):
                nonlocal replace_count
                original_text = run.text
                new_text = original_text
                
                for elem in sensitive_elements:
                    search_str = elem.get("text")
                    category = elem.get("category", "OTHER")
                    
                    if not search_str or not is_category_enabled(category, mask_options):
                        continue
                        
                    if re.search(re.escape(search_str), new_text, re.IGNORECASE):
                        replacement = get_text_replacement(search_str, category, mask_style)
                        pattern = re.compile(re.escape(search_str), re.IGNORECASE)
                        new_text = pattern.sub(replacement, new_text)
                        
                if new_text != original_text:
                    run.text = new_text
                    replace_count += 1
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            for run in p.runs:
                                replace_text_in_run(run)
                                
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                for p in cell.text_frame.paragraphs:
                                    for run in p.runs:
                                        replace_text_in_run(run)
                                        
            prs.save(output_path)
            logs.append({"type": "agent", "message": f"PPTX masking complete. Saved presentation. Applied {replace_count} run modifications."})
            return logs
        except Exception as e:
            logs.append({"type": "error", "message": f"Failed to mask PPTX: {str(e)}"})
            traceback.print_exc()
            return logs

    def mask_docx_file(self, input_path, output_path, mask_options, custom_keywords=None, custom_regex=None, user_prompt=None):
        """
        Masks a Word Document (.docx).
        """
        logs = [{"type": "agent", "message": "Starting DOCX masking process..."}]
        mask_style = mask_options.get("style", "redact")
        
        try:
            doc = Document(input_path)
            logs.append({"type": "info", "message": "Opened DOCX document."})
            
            all_text_runs = []
            for p in doc.paragraphs:
                for run in p.runs:
                    if run.text.strip():
                        all_text_runs.append(run.text)
                        
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for p in cell.paragraphs:
                            for run in p.runs:
                                if run.text.strip():
                                    all_text_runs.append(run.text)
                                    
            if not all_text_runs:
                logs.append({"type": "warning", "message": "No text detected in document paragraphs."})
                doc.save(output_path)
                return logs
                
            full_text_block = "\n---\n".join(all_text_runs)
            data, gemini_logs = self.analyze_text_for_pii(full_text_block, custom_keywords, custom_regex, mask_options, user_prompt)
            logs.extend(gemini_logs)
            
            sensitive_elements = data.get("sensitive_elements", [])
            if not sensitive_elements:
                logs.append({"type": "info", "message": "No sensitive elements identified in document."})
                doc.save(output_path)
                return logs
                
            replace_count = 0
            
            def replace_text_in_run(run):
                nonlocal replace_count
                original_text = run.text
                new_text = original_text
                
                for elem in sensitive_elements:
                    search_str = elem.get("text")
                    category = elem.get("category", "OTHER")
                    
                    if not search_str or not is_category_enabled(category, mask_options):
                        continue
                        
                    if re.search(re.escape(search_str), new_text, re.IGNORECASE):
                        replacement = get_text_replacement(search_str, category, mask_style)
                        pattern = re.compile(re.escape(search_str), re.IGNORECASE)
                        new_text = pattern.sub(replacement, new_text)
                        
                if new_text != original_text:
                    run.text = new_text
                    replace_count += 1
            
            for p in doc.paragraphs:
                for run in p.runs:
                    replace_text_in_run(run)
                    
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for p in cell.paragraphs:
                            for run in p.runs:
                                replace_text_in_run(run)
                                
            doc.save(output_path)
            logs.append({"type": "agent", "message": f"DOCX masking complete. Saved document. Applied {replace_count} run modifications."})
            return logs
        except Exception as e:
            logs.append({"type": "error", "message": f"Failed to mask DOCX: {str(e)}"})
            traceback.print_exc()
            return logs

    def mask_text_file(self, input_path, output_path, mask_options, custom_keywords=None, custom_regex=None, user_prompt=None):
        """
        Masks a plain text file (.txt, .csv, .json, etc.).
        """
        logs = [{"type": "agent", "message": "Starting plain text masking process..."}]
        mask_style = mask_options.get("style", "redact")
        
        try:
            with open(input_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
                
            logs.append({"type": "info", "message": f"Opened text file ({len(content)} characters)."})
            
            data, gemini_logs = self.analyze_text_for_pii(content, custom_keywords, custom_regex, mask_options, user_prompt)
            logs.extend(gemini_logs)
            
            sensitive_elements = data.get("sensitive_elements", [])
            if not sensitive_elements:
                logs.append({"type": "info", "message": "No sensitive elements identified."})
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(content)
                return logs
                
            masked_content = content
            replace_count = 0
            
            sensitive_elements.sort(key=lambda x: len(x.get("text", "")), reverse=True)
            
            for elem in sensitive_elements:
                search_str = elem.get("text")
                category = elem.get("category", "OTHER")
                
                if not search_str or not is_category_enabled(category, mask_options):
                    continue
                    
                if re.search(re.escape(search_str), masked_content, re.IGNORECASE):
                    replacement = get_text_replacement(search_str, category, mask_style)
                    pattern = re.compile(re.escape(search_str), re.IGNORECASE)
                    masked_content = pattern.sub(replacement, masked_content)
                    replace_count += 1
                    logs.append({"type": "success", "message": f"Redacted '{search_str}' -> '{replacement}'"})
                    
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(masked_content)
                
            logs.append({"type": "agent", "message": f"Text file masking complete. Replaced {replace_count} terms."})
            return logs
        except Exception as e:
            logs.append({"type": "error", "message": f"Failed to mask text: {str(e)}"})
            traceback.print_exc()
            return logs
